"""Generate the MemeBro Sprint Status Video 1 deck as a branded .pptx.

Run from the repo root (or anywhere):
    python admin/videos/generate_deck.py

Requires:
    pip install python-pptx

Output: admin/videos/status-video-1.pptx

The deck mirrors the MemeBro brand defined in styles/tokens.css:
warm cream background, orange accent, dark brown ink, Geist + Instrument Serif
fonts, rounded surface cards with a soft warm drop shadow. The actual mockups
from Design/screens/ are embedded directly into the visual slides.

If Geist or Instrument Serif aren't installed system-wide on the machine that
opens the deck, PowerPoint will substitute defaults. Install them from:
    https://vercel.com/font  (Geist)
    https://fonts.google.com/specimen/Instrument+Serif
"""

from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt


# ---------- Brand tokens (mirror styles/tokens.css exactly) ----------
CREAM     = RGBColor(0xF6, 0xEF, 0xDE)
SURFACE   = RGBColor(0xFC, 0xF8, 0xEE)
SURFACE_2 = RGBColor(0xFB, 0xF3, 0xDF)
LINE      = RGBColor(0xE6, 0xD6, 0xB0)
LINE_STR  = RGBColor(0xD4, 0xBE, 0x8E)
INK       = RGBColor(0x1F, 0x14, 0x08)
INK_2     = RGBColor(0x4A, 0x3A, 0x26)
INK_3     = RGBColor(0x8A, 0x79, 0x58)
ORANGE    = RGBColor(0xE3, 0x69, 0x1A)
ORANGE_2  = RGBColor(0xF0, 0x88, 0x38)
ORANGE_S  = RGBColor(0xFA, 0xD8, 0xB4)
ORANGE_W  = RGBColor(0xFC, 0xE6, 0xC8)
ORANGE_D  = RGBColor(0xB0, 0x4C, 0x0C)
MINT      = RGBColor(0x6F, 0xB8, 0x9C)
ROSE      = RGBColor(0xD9, 0x73, 0x84)
BLUE      = RGBColor(0x6A, 0x8F, 0xB8)

SANS  = "Geist"
SERIF = "Instrument Serif"

HERE       = Path(__file__).resolve().parent
REPO_ROOT  = HERE.parent.parent
DESIGN_DIR = REPO_ROOT / "Design"
OUTPUT     = HERE / "status-video-1.pptx"

FOOTER       = "MemeBro   ·   CSE 110 Group 16   ·   Sprint Status Video 1"
TOTAL_SLIDES = 8


# ---------- Low-level helpers ----------

def set_background(slide, color):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def add_warm_shadow(shape, *, blur=80000, dist=30000, alpha=22000):
    """No-op. Custom shadow XML broke Keynote compatibility; the cream-vs-surface
    fill contrast carries the card visual hierarchy on its own."""
    return


def add_textbox(slide, left, top, width, height, text, *,
                font=SANS, size=14, color=INK, bold=False, italic=False,
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                line_spacing=1.2):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return box


def add_card(slide, left, top, width, height, *,
             fill=SURFACE, line_color=None, radius_pct=8, shadow=True):
    """Rounded surface card. radius_pct is corner radius as 0..50."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, top, width, height)
    shape.adjustments[0] = radius_pct / 100
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.75)
    if shadow:
        add_warm_shadow(shape)
    return shape


def add_bullets(slide, left, top, width, height, items, *,
                size=15, color=INK, bullet_color=ORANGE,
                line_spacing=1.4, space_after_pt=10):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after_pt)
        dot = p.add_run()
        dot.text = "●  "
        dot.font.name = SANS
        dot.font.size = Pt(max(size - 3, 9))
        dot.font.color.rgb = bullet_color
        body = p.add_run()
        body.text = item
        body.font.name = SANS
        body.font.size = Pt(size)
        body.font.color.rgb = color
    return box


def add_image_card(slide, image_path, left, top, width, height, *,
                   caption=None, pad_in=0.18):
    """Embed an image inside a rounded surface card with a warm shadow."""
    add_card(slide, left, top, width, height,
             fill=SURFACE, shadow=True, radius_pct=6)
    if Path(image_path).exists():
        pad = Inches(pad_in)
        cap_reserve = Inches(0.35) if caption else Emu(0)
        slide.shapes.add_picture(
            str(image_path),
            left + pad, top + pad,
            width=width - 2 * pad,
            height=height - 2 * pad - cap_reserve,
        )
    if caption:
        add_textbox(slide, left, top + height - Inches(0.35),
                    width, Inches(0.3), caption,
                    font=SANS, size=10, color=INK_3, align=PP_ALIGN.CENTER)


def add_chrome(slide, num):
    """Footer + slide number, every non-title slide."""
    add_textbox(slide, Inches(0.7), Inches(7.05),
                Inches(10), Inches(0.35), FOOTER,
                size=9, color=INK_3)
    add_textbox(slide, Inches(11.5), Inches(7.05),
                Inches(1.2), Inches(0.35),
                f"{num}  /  {TOTAL_SLIDES}",
                size=9, color=INK_3, align=PP_ALIGN.RIGHT)


def add_title(slide, text):
    """Slide title in serif italic with an orange underline."""
    add_textbox(slide, Inches(0.7), Inches(0.45),
                Inches(12), Inches(0.9), text,
                font=SERIF, size=34, color=INK, italic=True)
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.7), Inches(1.28),
                                  Inches(0.6), Inches(0.06))
    rule.fill.solid()
    rule.fill.fore_color.rgb = ORANGE
    rule.line.fill.background()


def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, CREAM)
    return slide


# ---------- Slide builders ----------

def build_slide_1(prs):
    """Title slide."""
    slide = blank_slide(prs)
    # "MemeBro" big serif italic, centered
    add_textbox(slide, Inches(0), Inches(2.4),
                Inches(13.333), Inches(2.0),
                "MemeBro", font=SERIF, size=110, color=INK, italic=True,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Orange rule underneath, centered
    rule_w = Inches(2.5)
    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(13.333 / 2) - rule_w / 2,
        Inches(4.45), rule_w, Inches(0.08),
    )
    rule.fill.solid()
    rule.fill.fore_color.rgb = ORANGE
    rule.line.fill.background()
    # Subtitle stack
    add_textbox(slide, Inches(0), Inches(4.85),
                Inches(13.333), Inches(0.5),
                "Sprint Status Video 1",
                font=SANS, size=22, color=INK_2, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0), Inches(5.35),
                Inches(13.333), Inches(0.4),
                "CSE 110   ·   Spring 2026   ·   Group 16",
                font=SANS, size=13, color=INK_3, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0), Inches(5.75),
                Inches(13.333), Inches(0.4),
                "2026-05-20",
                font=SANS, size=11, color=INK_3, align=PP_ALIGN.CENTER)
    set_notes(slide,
        "Hi, we're group 16 and this is our mid-quarter status update on "
        "MemeBro, our AI-assisted meme generator. In the next four minutes "
        "we'll cover what we've built, how our team actually works, what we "
        "got wrong along the way, and what's coming next."
    )


def build_slide_2(prs):
    """What we're building."""
    slide = blank_slide(prs)
    add_title(slide, "What we're building")
    add_bullets(
        slide, Inches(0.7), Inches(2.1),
        Inches(6.0), Inches(4.5),
        [
            "Pick a meme template",
            "Type a vibe, AI suggests a caption",
            "Save it to your personal gallery",
            "Mobile first, runs entirely in the browser",
        ],
        size=18, line_spacing=1.5, space_after_pt=14,
    )
    img = DESIGN_DIR / "screens" / "01-library-home" / "1-mobile.png"
    add_image_card(slide, img,
                   Inches(7.6), Inches(1.7),
                   Inches(4.6), Inches(5.0),
                   caption="Library Home, mobile")
    add_chrome(slide, 2)
    set_notes(slide,
        "MemeBro is a web app where you pick a meme template, describe the "
        "vibe you want, and our AI suggests a caption. You can save what you "
        "make in a personal gallery. It's mobile first, no install, runs "
        "entirely in the browser today. Our target user is anyone who wants "
        "to make a quick joke without learning a design tool."
    )


def build_slide_3(prs):
    """Where we are."""
    slide = blank_slide(prs)
    add_title(slide, "Where we are")
    add_bullets(
        slide, Inches(0.7), Inches(2.1),
        Inches(7.5), Inches(4.5),
        [
            "5 active PRs across all 4 lanes, all in review",
            "Gallery, conjure UI, home page, ImgFlip fetch in flight",
            "Design system tokens and AGENTS.md already merged",
            "3 ADRs published: stack, hosting, backend strategy",
            "Deployment is sprint 3 day 1 (Cloudflare Pages setup in flight)",
        ],
        size=16, line_spacing=1.5, space_after_pt=10,
    )
    thumbs = [
        (DESIGN_DIR / "screens" / "01-library-home" / "1-mobile.png", "Library"),
        (DESIGN_DIR / "screens" / "03-conjure"      / "1-mobile.png", "Conjure"),
        (DESIGN_DIR / "screens" / "05-my-memes"     / "1-mobile.png", "My Memes"),
    ]
    thumb_w = Inches(1.35)
    thumb_h = Inches(2.7)
    gap     = Inches(0.18)
    start_x = Inches(8.45)
    for i, (path, label) in enumerate(thumbs):
        x = start_x + i * (thumb_w + gap)
        add_image_card(slide, path, x, Inches(2.3),
                       thumb_w, thumb_h, caption=label, pad_in=0.12)
    add_chrome(slide, 3)
    set_notes(slide,
        "We're at the walking skeleton stage. All four of our lanes have work "
        "in flight, with five active pull requests in review right now: the "
        "template gallery component, the conjure page UI, the home page "
        "layout, the ImgFlip API integration, and the CI/CD pipeline. We've "
        "also published three architecture decision records covering our "
        "stack, hosting, and backend strategy. Deployment itself is the first "
        "thing we land in sprint 3, on Cloudflare Pages."
    )


def build_slide_4(prs):
    """CI/CD pipeline."""
    slide = blank_slide(prs)
    add_title(slide, "CI/CD pipeline")
    add_bullets(
        slide, Inches(0.7), Inches(2.1),
        Inches(12), Inches(2.5),
        [
            "Lint, unit tests, end-to-end on every push and every PR",
            "Nothing merges to main without green checks",
            "Branch protection on main, PR only path",
            "Deploy target: Cloudflare Pages (sprint 3)",
            "Pipeline in final review, lands this week",
        ],
        size=15, line_spacing=1.4, space_after_pt=6,
    )
    # Horizontal flowchart at the bottom
    nodes = ["Push", "Lint", "Unit", "E2E", "Deploy"]
    n = len(nodes)
    node_w = Inches(1.7)
    node_h = Inches(0.95)
    chart_y = Inches(5.3)
    total_w = Inches(12)
    avail = total_w - node_w * n
    gap = avail / (n - 1)
    start_x = Inches(0.7)
    for i, label in enumerate(nodes):
        x = start_x + i * (node_w + gap)
        add_card(slide, x, chart_y, node_w, node_h,
                 fill=SURFACE, shadow=True, radius_pct=20)
        add_textbox(slide, x, chart_y, node_w, node_h, label,
                    font=SANS, size=15, color=INK, bold=True,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if i < n - 1:
            arrow_x = x + node_w
            add_textbox(slide, arrow_x, chart_y, gap, node_h, "→",
                        font=SANS, size=26, color=ORANGE, bold=True,
                        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_chrome(slide, 4)
    set_notes(slide,
        "Our CI/CD pipeline is built and in final review this week. Once "
        "merged, every push and every pull request runs lint, unit tests, "
        "and end-to-end tests. Nothing merges to main without all gates "
        "green, and main is already protected so the only path in is a "
        "reviewed pull request. We pivoted our hosting target from GitHub "
        "Pages to Cloudflare Pages this sprint because Pages Functions "
        "gives us a clean path to the AI proxy backend without changing "
        "how we deploy."
    )


def build_slide_5(prs):
    """How we actually work."""
    slide = blank_slide(prs)
    add_title(slide, "How we actually work")
    add_bullets(
        slide, Inches(0.7), Inches(2.1),
        Inches(7.0), Inches(4.7),
        [
            "Project lead + 4 lanes of 2 (frontend, backend, testing/docs, design)",
            "Standups at least 3 per week, all logged in the repo",
            "Weekly TA meeting with action items tracked",
            "All decisions land in ADRs, not Slack",
            "AGENTS.md keeps human and AI code consistent",
            "GitHub Issues are the source of truth, not Slack",
        ],
        size=14, line_spacing=1.45, space_after_pt=8,
    )
    # Mini org chart on the right: Lead + 4 lanes in a 2x2 grid
    chart_left = Inches(8.1)
    chart_top  = Inches(2.2)
    chart_w    = Inches(4.4)
    # Lead card centered above the grid
    lead_w = Inches(1.8)
    lead_x = chart_left + (chart_w - lead_w) / 2
    add_card(slide, lead_x, chart_top, lead_w, Inches(0.7),
             fill=ORANGE_W, shadow=True, radius_pct=20)
    add_textbox(slide, lead_x, chart_top, lead_w, Inches(0.7),
                "Project Lead", font=SANS, size=12, color=ORANGE_D, bold=True,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Connector line down to grid
    conn = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        chart_left + chart_w / 2 - Inches(0.02),
        chart_top + Inches(0.7),
        Inches(0.04), Inches(0.3),
    )
    conn.fill.solid()
    conn.fill.fore_color.rgb = LINE_STR
    conn.line.fill.background()
    # 2x2 grid of lane cards with a colored stripe each
    lanes = [
        ("Frontend",       MINT),
        ("Backend",        BLUE),
        ("Testing & Docs", ROSE),
        ("Design",         ORANGE_2),
    ]
    lane_w  = Inches(2.05)
    lane_h  = Inches(0.85)
    lane_gp = Inches(0.15)
    grid_top = chart_top + Inches(1.1)
    for i, (name, color) in enumerate(lanes):
        row, col = divmod(i, 2)
        x = chart_left + col * (lane_w + lane_gp)
        y = grid_top   + row * (lane_h + lane_gp)
        add_card(slide, x, y, lane_w, lane_h,
                 fill=SURFACE, shadow=True, radius_pct=15)
        stripe = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, y, Inches(0.1), lane_h)
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = color
        stripe.line.fill.background()
        add_textbox(slide, x + Inches(0.25), y,
                    lane_w - Inches(0.3), lane_h,
                    name, font=SANS, size=12, color=INK, bold=True,
                    anchor=MSO_ANCHOR.MIDDLE)
    add_chrome(slide, 5)
    set_notes(slide,
        "Process is what makes a ten person team actually ship. We're "
        "organized as a project lead plus four lanes of two for the build "
        "sprints: frontend, backend, testing and docs, and design. Standups "
        "happen at least three times a week and every one of them is logged "
        "in our repo. We meet with our TA every week with formal action items "
        "captured each session. All architectural decisions get written up "
        "as ADRs instead of being talked through in Slack. We also wrote an "
        "AGENTS.md file that defines our coding norms so humans and AI agents "
        "writing code stay consistent across the codebase. Our issue tracker "
        "is the source of truth and Slack is just where we coordinate."
    )


def build_slide_6(prs):
    """Challenges, and how we navigated them."""
    slide = blank_slide(prs)
    add_title(slide, "Challenges, and how we navigated them")
    rows = [
        ("Backend stack decision was non-trivial",
         "Structured deferral ADR with a re-decision date"),
        ("Hosting target evolved with AI integration shape",
         "Pivoted to Cloudflare Pages, captured in an ADR"),
        ("Testing infrastructure was bigger than one issue",
         "Split into lint, unit, and end-to-end sub-issues"),
        ("Coordinating AI-assisted code across many contributors",
         "AGENTS.md as a shared contract for the codebase"),
        ("Cross-team PR review at scale",
         "24 hour review SLA + lane Slack channels"),
    ]
    table_left   = Inches(0.7)
    table_top    = Inches(2.0)
    challenge_w  = Inches(5.6)
    sol_gap      = Inches(0.4)
    sol_w        = Inches(12) - challenge_w - sol_gap
    sol_x        = table_left + challenge_w + sol_gap
    # Headers
    add_textbox(slide, table_left, table_top, challenge_w, Inches(0.4),
                "Challenge", font=SANS, size=11, color=INK_3, bold=True)
    add_textbox(slide, sol_x, table_top, sol_w, Inches(0.4),
                "How we navigated it",
                font=SANS, size=11, color=ORANGE, bold=True)
    # Hairline rule under headers
    hr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                table_left, table_top + Inches(0.45),
                                Inches(12), Inches(0.015))
    hr.fill.solid()
    hr.fill.fore_color.rgb = LINE
    hr.line.fill.background()
    row_h = Inches(0.78)
    for i, (challenge, solution) in enumerate(rows):
        y = table_top + Inches(0.6) + i * (row_h + Inches(0.1))
        add_textbox(slide, table_left + Inches(0.1), y,
                    challenge_w - Inches(0.1), row_h,
                    challenge, font=SANS, size=13, color=INK_2,
                    anchor=MSO_ANCHOR.MIDDLE)
        add_card(slide, sol_x, y, sol_w, row_h,
                 fill=ORANGE_W, shadow=False, radius_pct=20)
        add_textbox(slide, sol_x + Inches(0.3), y,
                    sol_w - Inches(0.5), row_h,
                    solution, font=SANS, size=13, color=ORANGE_D,
                    anchor=MSO_ANCHOR.MIDDLE)
    add_chrome(slide, 6)
    set_notes(slide,
        "Every team hits real friction, and ours has been around making hard "
        "decisions cleanly. The backend stack was a non-trivial choice "
        "between static-only, serverless functions, and a full backend, so "
        "instead of picking one prematurely we wrote a structured deferral "
        "ADR with a clear re-decision date. As our AI integration came into "
        "focus, we realized our original hosting target wasn't the right fit "
        "and pivoted to Cloudflare Pages, with the rationale captured in an "
        "ADR. Our testing infrastructure work was larger than a single issue, "
        "so we split it into lint, unit, and end-to-end sub-items that can "
        "ship in parallel. With multiple contributors using AI assistants on "
        "the same repo, we wrote AGENTS.md as a shared contract to keep "
        "style and conventions coherent. And we set a 24 hour review SLA "
        "with lane-specific Slack channels so pull requests don't sit in "
        "queue."
    )


def build_slide_7(prs):
    """AI usage."""
    slide = blank_slide(prs)
    add_title(slide, "AI usage")
    columns = [
        ("Pros", MINT, [
            "Faster scaffolding (AGENTS.md, tokens, ADR drafts)",
            "Cleaner JSDoc coverage across the codebase",
            "Meeting transcripts to structured notes",
            "Agents follow our written conventions",
        ]),
        ("Cons", ROSE, [
            "Will invent abstractions if not constrained",
            "Output needs human review before commit",
            "Easy to over-document and create noise",
        ]),
        ("How we mitigate", ORANGE, [
            "AGENTS.md acts as a contract for AI agents",
            "JSDoc lint enforces consistency",
            "Human author owns and reviews every PR",
        ]),
    ]
    col_top = Inches(2.0)
    col_h   = Inches(4.6)
    col_gap = Inches(0.3)
    total_w = Inches(12)
    col_w   = (total_w - col_gap * 2) / 3
    start_x = Inches(0.7)
    for i, (name, accent, items) in enumerate(columns):
        x = start_x + i * (col_w + col_gap)
        add_card(slide, x, col_top, col_w, col_h,
                 fill=SURFACE, shadow=True, radius_pct=8)
        # Accent stripe along the top
        stripe = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x, col_top, col_w, Inches(0.12))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = accent
        stripe.line.fill.background()
        # Header inside the card
        add_textbox(slide, x + Inches(0.35), col_top + Inches(0.3),
                    col_w - Inches(0.7), Inches(0.5),
                    name, font=SANS, size=18, color=INK, bold=True)
        # Bullets inside the card
        add_bullets(slide, x + Inches(0.35), col_top + Inches(0.95),
                    col_w - Inches(0.7), col_h - Inches(1.1),
                    items, size=12, line_spacing=1.4, space_after_pt=8,
                    bullet_color=accent)
    add_chrome(slide, 7)
    set_notes(slide,
        "AI has been a real force multiplier. We use it for scaffolding, "
        "drafting ADRs, transcribing meeting notes, and keeping our "
        "documentation in shape. The cons are real though. Agents will "
        "invent abstractions if you let them, output needs human review "
        "before merge, and it's very easy to over-document. We mitigate by "
        "treating AGENTS.md as a contract that bounds what AI does, by "
        "enforcing JSDoc through lint, and by requiring the human author to "
        "own and review every PR."
    )


def build_slide_8(prs):
    """What's next."""
    slide = blank_slide(prs)
    add_title(slide, "What's next")
    add_bullets(
        slide, Inches(0.7), Inches(2.0),
        Inches(12), Inches(2.0),
        [
            "Sprint 3: Cloudflare Pages deploy live, AI proxy, end-to-end flow, test coverage",
            "Sprint 4: polish, accessibility pass, final presentation prep",
            "Working end-to-end demo by sprint 3 close",
            "Full product walkthrough at the final presentation",
        ],
        size=14, line_spacing=1.4, space_after_pt=8,
    )
    # Timeline at the bottom
    nodes = [
        ("Sprint 3", "Deploy + AI proxy",  ORANGE),
        ("Sprint 4", "Polish + a11y",      MINT),
        ("Final",    "Product walkthrough", BLUE),
    ]
    n        = len(nodes)
    node_w   = Inches(2.6)
    line_y   = Inches(5.55)
    total_w  = Inches(12)
    node_total = node_w * n
    gap      = (total_w - node_total) / (n - 1) if n > 1 else Inches(0)
    start_x  = Inches(0.7)
    # Horizontal track
    track = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        start_x + Inches(0.3), line_y,
        total_w - Inches(0.6), Inches(0.04),
    )
    track.fill.solid()
    track.fill.fore_color.rgb = LINE_STR
    track.line.fill.background()
    for i, (label, sub, color) in enumerate(nodes):
        x = start_x + i * (node_w + gap)
        # Label above the dot
        add_textbox(slide, x, line_y - Inches(0.55),
                    node_w, Inches(0.45),
                    label, font=SANS, size=15, color=INK, bold=True,
                    align=PP_ALIGN.CENTER)
        # Dot on the track
        dot_size = Inches(0.32)
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            x + node_w / 2 - dot_size / 2,
            line_y - dot_size / 2 + Inches(0.02),
            dot_size, dot_size,
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.color.rgb = CREAM
        dot.line.width = Pt(3)
        # Sub below the dot
        add_textbox(slide, x, line_y + Inches(0.35),
                    node_w, Inches(0.45),
                    sub, font=SANS, size=11, color=INK_2,
                    align=PP_ALIGN.CENTER)
    # Closing line
    add_textbox(slide, Inches(0), Inches(6.55),
                Inches(13.333), Inches(0.45),
                "Thanks, Group 16.",
                font=SERIF, size=20, color=ORANGE, italic=True,
                align=PP_ALIGN.CENTER)
    add_chrome(slide, 8)
    set_notes(slide,
        "Sprint 3 starts Sunday. Day one we land the Cloudflare Pages deploy "
        "so we have a live URL. From there we stand up the AI proxy on Pages "
        "Functions, wire the full end-to-end meme flow, and bring test "
        "coverage up across the components we've already shipped. Sprint 4 "
        "is polish, an accessibility pass, and final presentation prep. By "
        "the close of sprint 3 we'll have a working end-to-end demo, and "
        "we'll bring a full product walkthrough to the final presentation. "
        "Thanks for watching."
    )


# ---------- Main ----------

def _fix_pptx_metadata(prs):
    """Fix slide-size type attribute and app.xml counters so the file opens
    cleanly in PowerPoint, Keynote, and Google Slides."""
    sld_sz = prs.element.find(qn("p:sldSz"))
    if sld_sz is not None:
        sld_sz.set("type", "custom")

    app_part = prs.part.package.part_related_by(
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/extended-properties"
    )
    app_el = etree.fromstring(app_part.blob)
    ns = {"ep": "http://schemas.openxmlformats.org/officeDocument/2006/extended-properties"}
    slides_el = app_el.find("ep:Slides", ns)
    if slides_el is not None:
        slides_el.text = str(len(prs.slides))
    fmt_el = app_el.find("ep:PresentationFormat", ns)
    if fmt_el is not None:
        fmt_el.text = "Widescreen"
    app_part._blob = etree.tostring(app_el, xml_declaration=True,
                                    encoding="UTF-8", standalone=True)


def main():
    prs = Presentation()
    prs.slide_width  = Emu(12192000)
    prs.slide_height = Emu(6858000)

    build_slide_1(prs)
    build_slide_2(prs)
    build_slide_3(prs)
    build_slide_4(prs)
    build_slide_5(prs)
    build_slide_6(prs)
    build_slide_7(prs)
    build_slide_8(prs)

    _fix_pptx_metadata(prs)
    prs.save(OUTPUT)
    print(f"Wrote {OUTPUT}")
    print()
    print("Fonts: this deck uses Geist and Instrument Serif.")
    print("If those aren't installed system-wide, PowerPoint substitutes a")
    print("default sans/serif. Install them for the on-brand rendering:")
    print("  Geist:            https://vercel.com/font")
    print("  Instrument Serif: https://fonts.google.com/specimen/Instrument+Serif")


if __name__ == "__main__":
    main()
